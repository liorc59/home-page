"""
MiAmar Pre-Seed Pitch Deck Generator
Generates a 14-slide VC pitch deck as .pptx using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──────────────────────────────────────────────
BLUE = RGBColor(0x1E, 0x40, 0xAF)
TEAL = RGBColor(0x0D, 0x94, 0x88)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_TEAL_BG = RGBColor(0xF0, 0xFD, 0xFA)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED_ACCENT = RGBColor(0xDC, 0x26, 0x26)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bullet_text(tf, text, font_size=16, color=DARK, bold=False, indent_level=0,
                    space_after=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_BODY
    p.level = indent_level
    p.space_after = space_after
    return p


def add_multi_line_textbox(slide, left, top, width, height, lines, default_size=16,
                           default_color=DARK, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Blue accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_textbox(slide, Inches(2), Inches(1.8), Inches(9.333), Inches(1.2),
                "MiAmar", font_size=60, color=BLUE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_textbox(slide, Inches(2), Inches(3.1), Inches(9.333), Inches(0.8),
                "AI-Powered Competitive Intelligence for Product Teams",
                font_size=28, color=DARK, bold=False, alignment=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.5), Inches(4.2), Inches(2.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_textbox(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
                "Pre-Seed Round — $1M", font_size=22, color=GRAY,
                bold=False, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "Thank you for your time. I'm Eran Menachemi, CEO of MiAmar. We're building the competitive intelligence platform that product marketing and product teams have been asking for — one that actually works without a dedicated analyst team. We're raising a $1M pre-seed round to capture a market that's being underserved by expensive, manual-heavy incumbents.")


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "PMMs Are Drowning in Manual Competitive Intel",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # Hero stat box
    hero = add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.3),
                            LIGHT_BLUE_BG, BLUE)
    hero_tb = add_textbox(slide, Inches(1.2), Inches(1.65), Inches(11), Inches(1.0),
                          "PMMs spend 6-8 hours/week on manual competitor monitoring — and still miss critical signals.",
                          font_size=24, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Stats grid (2x2)
    stats = [
        ("90% of businesses", "say their industry has become\nmore competitive in the last 3 years"),
        ("65% of opportunities", "are competitive — yet average team\nrates itself 3.8/10 in CI effectiveness"),
        ("CI teams shrinking", "Companies with 3+ CI staff dropped\nfrom 34% to 25% in two years"),
        ("$2M–$10M/year", "in deals companies could have won\n— the cost of the CI gap"),
    ]

    x_positions = [Inches(0.8), Inches(6.9)]
    y_positions = [Inches(3.2), Inches(5.2)]

    for idx, (title, desc) in enumerate(stats):
        col = idx % 2
        row = idx // 2
        x = x_positions[col]
        y = y_positions[row]
        w = Inches(5.8)
        h = Inches(1.7)

        card = add_rounded_rect(slide, x, y, w, h, LIGHT_BG)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.5),
                    title, font_size=20, color=BLUE, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.75), w - Inches(0.6), Inches(0.8),
                    desc, font_size=14, color=GRAY)

    set_notes(slide,
              "The competitive intelligence problem is simple: markets are getting more competitive, but the teams responsible for tracking competitors are getting smaller. PMMs — the people who own battlecards, positioning, and sales enablement — are spending 6 to 8 hours every week manually scanning LinkedIn, news, and job boards. And they're still missing things. Crayon's research shows this gap costs companies $2 to $10 million per year in winnable deals. The current tools haven't solved this — they've just moved the manual work from spreadsheets to another platform.")


def slide_03_market_pain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "This Is What PMMs Are Saying Right Now",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4),
                "Direct quotes from Reddit's r/ProductMarketing and r/ProductManagement (27 posts analyzed)",
                font_size=14, color=GRAY)

    quotes = [
        ('"My director wants me to spend 4 hours/day monitoring competitor news — but only gathering, no analysis. Pure data collection drudgery."',
         "PMM, r/ProductMarketing", "10/10"),
        ('"I tried Klue but found the data goes stale fast — it\'s basically a storage system that still requires you to manually feed it intel."',
         "PM, r/ProductManagement", "10/10"),
        ('"SMB looking for CI tools — Klue and Crayon are $20K–$40K/year which makes no sense for us."',
         "PMM, r/ProductMarketing", "9/10"),
        ('"I tried ChatGPT for competitive analysis — output was generic and often wrong. NotebookLM was better but required tedious manual data entry."',
         "PMM, r/ProductMarketing", "8/10"),
    ]

    x_positions = [Inches(0.8), Inches(6.9)]
    y_positions = [Inches(1.8), Inches(4.5)]

    for idx, (quote, source, urgency) in enumerate(quotes):
        col = idx % 2
        row = idx // 2
        x = x_positions[col]
        y = y_positions[row]
        w = Inches(5.8)
        h = Inches(2.4)

        card = add_rounded_rect(slide, x, y, w, h, LIGHT_BG, GRAY)

        # Quote mark
        add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(0.5), Inches(0.5),
                    "\u201C", font_size=36, color=BLUE, bold=True)

        add_textbox(slide, x + Inches(0.3), y + Inches(0.4), w - Inches(0.6), Inches(1.3),
                    quote, font_size=13, color=DARK, bold=False)

        add_textbox(slide, x + Inches(0.3), y + Inches(1.7), w - Inches(1.8), Inches(0.3),
                    f"— {source}", font_size=12, color=GRAY, bold=False)

        # Urgency badge
        badge = add_rounded_rect(slide, x + w - Inches(1.4), y + Inches(1.75),
                                 Inches(1.1), Inches(0.35), ORANGE)
        add_textbox(slide, x + w - Inches(1.4), y + Inches(1.75),
                    Inches(1.1), Inches(0.35),
                    f"Urgency: {urgency}", font_size=10, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.35),
                "27 Reddit posts in the last 12 months actively asking for a tool like MiAmar.",
                font_size=13, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "These aren't hypothetical personas — these are real PMMs and PMs posting on Reddit in the last 12 months, actively searching for what we're building. The themes are consistent: manual monitoring is unsustainable, existing tools like Klue and Crayon are either too expensive or too manual, and generic AI tools like ChatGPT hallucinate when used for CI. We analyzed 27 posts — 19 showed explicit willingness to pay, and the average urgency rating was 8+ out of 10. This is a market that's actively looking for a solution and can't find one at the right price point.")


def slide_04_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "MiAmar: Competitive Intelligence on Autopilot",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # One-liner hero
    hero = add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
                            LIGHT_TEAL_BG, TEAL)
    add_textbox(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(0.8),
                "Add your competitors. AI monitors 24/7. Get actionable intelligence — not raw data.",
                font_size=22, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # 4 capability cards in a row
    capabilities = [
        ("Monitors", "LinkedIn posts, company profiles,\nnews, funding rounds, job\npostings — automatically, every day"),
        ("Classifies", "Every signal using AI: product\nlaunches, funding events, key\nhires, partnerships, positioning shifts"),
        ("Generates", "Ready-to-use battlecards,\ncompetitor profiles, prospect\ncards, weekly intelligence digests"),
        ("Delivers", "Insights where your team works:\nSlack, HubSpot, Monday.com,\nemail"),
    ]

    card_w = Inches(2.8)
    gap = Inches(0.2)
    start_x = Inches(0.55)

    for i, (title, desc) in enumerate(capabilities):
        x = start_x + i * (card_w + gap)
        y = Inches(3.0)
        h = Inches(2.8)

        card = add_rounded_rect(slide, x, y, card_w, h, LIGHT_BG, BLUE)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(1.0), y + Inches(0.2),
                                        Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        add_textbox(slide, x + Inches(1.0), y + Inches(0.25),
                    Inches(0.7), Inches(0.6),
                    str(i + 1), font_size=22, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.1), card_w - Inches(0.4), Inches(0.5),
                    title, font_size=20, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.6), card_w - Inches(0.4), Inches(1.1),
                    desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Differentiator
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
                "No manual curation. No stale data. No $40K/year price tag.",
                font_size=18, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "MiAmar is the competitive intelligence platform that PMMs have been asking for. You add your competitors — just paste a LinkedIn URL — and our AI starts monitoring them immediately. Every day, we scan LinkedIn posts, company profiles, news articles, funding rounds, and job postings. Our AI classifies every signal, generates battlecards and intelligence summaries, and delivers them where your team already works. The key difference: with Klue or Crayon, you're paying $20K to $40K per year for a platform that still requires manual input to stay current. With MiAmar, the intelligence generates itself.")


def slide_05_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Three Steps. Zero Manual Work.",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    steps = [
        ("1", "Add Companies",
         "Paste a LinkedIn URL or company name.\nMonitoring starts in seconds.",
         "Competitor tracked"),
        ("2", "AI Monitors 24/7",
         "Daily scans of LinkedIn, news, funding,\njob postings. AI classifies every signal.",
         "Continuous intelligence"),
        ("3", "Get Actionable Insights",
         "Auto-generated battlecards, event feeds,\nprospect cards, weekly newsletter.",
         "Ready-to-use intel"),
    ]

    card_w = Inches(3.6)
    gap = Inches(0.4)
    start_x = Inches(0.7)

    for i, (num, title, desc, output) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        y = Inches(1.8)
        h = Inches(4.5)

        card = add_rounded_rect(slide, x, y, card_w, h, LIGHT_BG, BLUE)

        # Big number
        num_box = add_rounded_rect(slide, x + Inches(1.25), y + Inches(0.3),
                                   Inches(1.1), Inches(1.1), BLUE)
        add_textbox(slide, x + Inches(1.25), y + Inches(0.35),
                    Inches(1.1), Inches(1.0),
                    num, font_size=44, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Inches(0.5),
                    title, font_size=20, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(2.2), card_w - Inches(0.6), Inches(1.0),
                    desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

        # Output badge
        badge = add_rounded_rect(slide, x + Inches(0.5), y + Inches(3.5),
                                 card_w - Inches(1.0), Inches(0.5), TEAL)
        add_textbox(slide, x + Inches(0.5), y + Inches(3.52),
                    card_w - Inches(1.0), Inches(0.45),
                    output, font_size=13, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Arrow between cards
        if i < 2:
            arrow_x = x + card_w + Inches(0.05)
            arrow_y = y + Inches(2.0)
            add_textbox(slide, arrow_x, arrow_y, Inches(0.3), Inches(0.5),
                        "\u2192", font_size=28, color=BLUE, bold=True,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "The workflow is dead simple. Step one: add the companies you want to track — paste a LinkedIn URL or type a name. Takes 10 seconds. Step two: our AI monitors those companies 24/7. We scan LinkedIn posts, company profiles, news, funding data, and job postings every single day and classify every signal automatically. Step three: you get ready-to-use output — auto-generated battlecards, an event intelligence feed, prospect cards for sales prep, and a weekly newsletter digest sent to your team's inbox. The PMM goes from spending 6 to 8 hours per week on monitoring to spending 15 minutes reviewing a curated digest.")


def slide_06_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "A Full CI Platform — Not Just a Dashboard",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    features = [
        ("Market Intelligence Dashboard",
         "AI-prioritized signal feed across all tracked companies. Trend detection."),
        ("Event Intelligence",
         "AI-classified events: product launches, funding, hires, partnerships."),
        ("Sales Battlecards",
         "Auto-generated, self-updating. Elevator pitch, objection handling, talk tracks."),
        ("AI Prospect & Client Cards",
         "Complete prospect snapshot. Pain points matched to your product fit."),
        ("Weekly CI Newsletter",
         "Automated intelligence digest delivered to your team's inbox."),
        ("Meeting Prep",
         "AI meeting briefs in 60 seconds. 19+ meeting types supported."),
    ]

    card_w = Inches(3.8)
    card_h = Inches(1.7)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(1.5)

    for idx, (title, desc) in enumerate(features):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_BG, BLUE)

        add_textbox(slide, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.5), Inches(0.5),
                    title, font_size=16, color=BLUE, bold=True)

        add_textbox(slide, x + Inches(0.25), y + Inches(0.8), card_w - Inches(0.5), Inches(0.7),
                    desc, font_size=13, color=GRAY)

    # Integration bar
    int_bar = add_rounded_rect(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.6),
                               LIGHT_TEAL_BG, TEAL)
    add_textbox(slide, Inches(0.7), Inches(5.42), Inches(11.9), Inches(0.55),
                "Integrations:   HubSpot   |   Slack   |   Monday.com",
                font_size=16, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "We're not building a single feature — we're building the full competitive intelligence workspace. The dashboard gives you a market-level view with AI-prioritized signals. Event intelligence catches and classifies every competitor move. Battlecards are auto-generated and self-updating — sales teams actually use them because they're always current. Prospect cards give your reps an unfair advantage before every meeting. The weekly newsletter keeps the entire team aligned without any manual effort. And meeting prep generates a complete brief in 60 seconds. All of this integrates with HubSpot, Slack, and Monday.com.")


def slide_07_market_size(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "The CI Tools Market Is Growing — and Underserved",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # TAM / SAM / SOM concentric circles via layered rounded rects
    cx = Inches(3.5)
    cy = Inches(1.6)

    # TAM - outermost
    tam = add_rounded_rect(slide, cx - Inches(0.2), cy, Inches(5.8), Inches(5.0),
                           LIGHT_BLUE_BG, BLUE)
    add_textbox(slide, cx, cy + Inches(0.15), Inches(5.4), Inches(0.4),
                "TAM: $50B+", font_size=22, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, cy + Inches(0.55), Inches(5.4), Inches(0.35),
                "Global competitive intelligence market",
                font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # SAM - middle
    sam = add_rounded_rect(slide, cx + Inches(0.6), cy + Inches(1.1),
                           Inches(4.2), Inches(3.2), LIGHT_TEAL_BG, TEAL)
    add_textbox(slide, cx + Inches(0.8), cy + Inches(1.3), Inches(3.8), Inches(0.4),
                "SAM: $500M", font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx + Inches(0.8), cy + Inches(1.7), Inches(3.8), Inches(0.35),
                "CI software tools → $1.5B by 2034 (12.7% CAGR)",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # SOM - innermost
    som = add_rounded_rect(slide, cx + Inches(1.4), cy + Inches(2.3),
                           Inches(2.6), Inches(1.5), WHITE, BLUE)
    add_textbox(slide, cx + Inches(1.5), cy + Inches(2.5), Inches(2.4), Inches(0.4),
                "SOM: $15M", font_size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx + Inches(1.5), cy + Inches(2.9), Inches(2.4), Inches(0.5),
                "~3,000 B2B SaaS companies\n× $5K ACV",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Supporting stats on the right
    right_x = Inches(9.0)
    stats_lines = [
        ("CI tools market", "12-13% CAGR growth"),
        ("AI SaaS market", "38% CAGR → $775B by 2031"),
        ("30K-72K", "B2B SaaS companies globally"),
    ]

    for i, (title, detail) in enumerate(stats_lines):
        y = Inches(2.2) + i * Inches(1.5)
        card = add_rounded_rect(slide, right_x, y, Inches(3.8), Inches(1.2), LIGHT_BG, GRAY)
        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.15),
                    Inches(3.4), Inches(0.4),
                    title, font_size=16, color=BLUE, bold=True)
        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.6),
                    Inches(3.4), Inches(0.4),
                    detail, font_size=13, color=GRAY)

    set_notes(slide,
              "The broad competitive intelligence market is over $50 billion, but that includes consulting and services. The relevant number for us is the CI tools market — approximately $500 million today, growing to $1.5 billion by 2034 at a 12 to 13 percent CAGR. Our initial serviceable market is mid-market B2B SaaS companies with a product marketing function — roughly 3,000 companies at our entry-level ACV, giving us a $15 million near-term opportunity. The AI SaaS tailwind is real: the broader AI SaaS market is growing at 38% annually. We're riding both the CI growth wave and the AI infrastructure cost collapse.")


def slide_08_why_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Three Tailwinds Making This Possible Today",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    tailwinds = [
        ("AI Cost Collapse",
         BLUE,
         [
             "AI infrastructure costs dropped by orders of magnitude",
             "15% cost savings, 23% productivity gains (Gartner)",
             "32% average cost savings for advanced adopters (Deloitte)",
             "What required analysts → now done by AI at fraction of cost",
         ]),
        ("Incumbent Pricing Gap",
         TEAL,
         [
             "Klue: $20K–$45K/yr. Crayon: $15K–$47K/yr (Vendr)",
             "Both still require manual curation to stay current",
             "SMBs and mid-market completely priced out",
             "Massive gap between $0 (manual) and $20K+ (enterprise)",
         ]),
        ("PMM Role Explosion",
         ORANGE,
         [
             "54,000+ PMM positions posted in 2025 (Zippia)",
             "Top companies: 1 PMM for every 1.6 PMs (McKinsey)",
             "53% of PMMs now measured directly on revenue",
             "More PMMs = more demand for CI at accessible prices",
         ]),
    ]

    card_w = Inches(3.8)
    gap = Inches(0.25)
    start_x = Inches(0.7)

    for i, (title, accent, bullets) in enumerate(tailwinds):
        x = start_x + i * (card_w + gap)
        y = Inches(1.5)
        h = Inches(5.3)

        card = add_rounded_rect(slide, x, y, card_w, h, LIGHT_BG)

        # Color bar at top of card
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y, card_w, Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_textbox(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.5),
                    title, font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

        # Bullets
        txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.0),
                                         card_w - Inches(0.6), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.font.name = FONT_BODY
            p.space_after = Pt(8)

    set_notes(slide,
              "Three things are converging right now. First, AI costs have collapsed — what used to require a dedicated analyst team can now be automated at a fraction of the cost. Second, the incumbents are charging $20K to $45K per year and still require manual work. There's a massive gap between free (spreadsheets and Google Alerts) and $20K+ (Klue/Crayon). That's where we play. Third, the PMM role is exploding. McKinsey found that the fastest-growing companies have significantly more PMMs per product manager, and over half of PMMs are now measured directly on revenue. More PMMs, more strategic pressure, more demand for affordable CI tools. The timing is perfect.")


def slide_09_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Land & Expand in Product-Led Teams",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # Pricing table
    tiers = [
        ("Starter", "$199/mo", "Solo PMMs, small teams",
         "Up to 10 competitors, core monitoring,\nbattlecards, weekly newsletter"),
        ("Growth", "$499/mo", "Mid-market PMM teams",
         "Up to 30 competitors, event intelligence,\nprospect cards, integrations, meeting prep"),
        ("Enterprise", "Custom", "Large PMM orgs",
         "Unlimited competitors, API access,\ndedicated support, custom integrations"),
    ]

    card_w = Inches(3.8)
    gap = Inches(0.25)
    start_x = Inches(0.7)

    for i, (name, price, target, features) in enumerate(tiers):
        x = start_x + i * (card_w + gap)
        y = Inches(1.4)
        h = Inches(3.0)

        is_highlight = (i == 1)
        bg_color = LIGHT_BLUE_BG if is_highlight else LIGHT_BG
        border = BLUE if is_highlight else GRAY

        card = add_rounded_rect(slide, x, y, card_w, h, bg_color, border)

        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.6), Inches(0.4),
                    name, font_size=20, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(0.65), card_w - Inches(0.6), Inches(0.5),
                    price, font_size=28, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(1.15), card_w - Inches(0.6), Inches(0.3),
                    target, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Inches(1.2),
                    features, font_size=12, color=DARK, alignment=PP_ALIGN.CENTER)

    # Target ACV
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.4),
                "Target ACV: $5K–$10K in Year 1, growing to $15K+ as teams expand",
                font_size=15, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # Land & Expand flow
    le_y = Inches(5.2)
    steps = [
        ("Land", "Solo PMM signs up, adds\n5–10 competitors.\nSees value in Week 1."),
        ("Expand", "Shares weekly newsletter\nwith sales team.\nSales asks for battlecards."),
        ("Upsell", "Team moves to Growth tier.\nAdds integrations.\nACV doubles."),
    ]

    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        h = Inches(1.8)

        accent = [BLUE, TEAL, ORANGE][i]
        card = add_rounded_rect(slide, x, le_y, card_w, h, WHITE, accent)

        # Top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, le_y, card_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_textbox(slide, x + Inches(0.2), le_y + Inches(0.2),
                    card_w - Inches(0.4), Inches(0.35),
                    f"{i + 1}. {title}", font_size=16, color=accent, bold=True)

        add_textbox(slide, x + Inches(0.2), le_y + Inches(0.6),
                    card_w - Inches(0.4), Inches(1.0),
                    desc, font_size=12, color=GRAY)

        if i < 2:
            arrow_x = x + card_w + Inches(0.02)
            add_textbox(slide, arrow_x, le_y + Inches(0.6), Inches(0.2), Inches(0.5),
                        "\u2192", font_size=24, color=BLUE, bold=True)

    set_notes(slide,
              "Our business model is designed for land and expand. A solo PMM signs up at $199 per month — accessible enough to start with a credit card. They add their top competitors and immediately get auto-generated battlecards and a weekly digest. Within a few weeks, they share the newsletter with the sales team. Sales asks for battlecards in Slack. The team upgrades to Growth. Our target ACV is $5K to $10K in Year 1, growing as teams expand. The pricing is deliberately set below the incumbent threshold — Klue and Crayon start at $15K to $20K. We can win the mid-market deals they're losing on price while still building a high-margin SaaS business.")


def slide_10_competitive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Positioned to Win the Mid-Market",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # 2x2 positioning matrix
    matrix_x = Inches(0.8)
    matrix_y = Inches(1.4)
    matrix_w = Inches(5.5)
    matrix_h = Inches(4.5)

    # Axis labels
    add_textbox(slide, matrix_x + Inches(1.0), matrix_y + matrix_h + Inches(0.05),
                Inches(3.5), Inches(0.3),
                "Manual Effort Required →",
                font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Y-axis label (vertical text is hard, use short text)
    add_textbox(slide, matrix_x - Inches(0.7), matrix_y + Inches(1.5),
                Inches(0.7), Inches(1.5),
                "↑\nDepth of\nIntelligence",
                font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Quadrant boxes
    half_w = matrix_w / 2
    half_h = matrix_h / 2

    # Top-left: MiAmar (best)
    q1 = add_rounded_rect(slide, matrix_x, matrix_y,
                          half_w - Inches(0.05), half_h - Inches(0.05),
                          LIGHT_TEAL_BG, TEAL)
    add_textbox(slide, matrix_x + Inches(0.15), matrix_y + Inches(0.3),
                half_w - Inches(0.4), Inches(0.4),
                "MiAmar", font_size=20, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, matrix_x + Inches(0.15), matrix_y + Inches(0.8),
                half_w - Inches(0.4), Inches(0.6),
                "AI-automated, full platform\nDeep + Low effort",
                font_size=12, color=DARK, alignment=PP_ALIGN.CENTER)

    # Top-right: Klue/Crayon
    q2 = add_rounded_rect(slide, matrix_x + half_w + Inches(0.05), matrix_y,
                          half_w - Inches(0.05), half_h - Inches(0.05),
                          LIGHT_BG, GRAY)
    add_textbox(slide, matrix_x + half_w + Inches(0.2), matrix_y + Inches(0.3),
                half_w - Inches(0.4), Inches(0.4),
                "Klue / Crayon", font_size=18, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, matrix_x + half_w + Inches(0.2), matrix_y + Inches(0.8),
                half_w - Inches(0.4), Inches(0.6),
                "$20K-$45K, requires analysts\nDeep + High effort",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom-left: Google Alerts + ChatGPT
    q3 = add_rounded_rect(slide, matrix_x, matrix_y + half_h + Inches(0.05),
                          half_w - Inches(0.05), half_h - Inches(0.05),
                          LIGHT_BG, GRAY)
    add_textbox(slide, matrix_x + Inches(0.15), matrix_y + half_h + Inches(0.3),
                half_w - Inches(0.4), Inches(0.4),
                "Google Alerts + ChatGPT", font_size=14, color=GRAY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, matrix_x + Inches(0.15), matrix_y + half_h + Inches(0.8),
                half_w - Inches(0.4), Inches(0.6),
                "Free but unreliable\nShallow + Low effort",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom-right: Manual research
    q4 = add_rounded_rect(slide, matrix_x + half_w + Inches(0.05), matrix_y + half_h + Inches(0.05),
                          half_w - Inches(0.05), half_h - Inches(0.05),
                          LIGHT_BG, GRAY)
    add_textbox(slide, matrix_x + half_w + Inches(0.2), matrix_y + half_h + Inches(0.3),
                half_w - Inches(0.4), Inches(0.4),
                "Manual Research", font_size=16, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, matrix_x + half_w + Inches(0.2), matrix_y + half_h + Inches(0.8),
                half_w - Inches(0.4), Inches(0.6),
                "Spreadsheets, 6-8 hrs/week\nShallow + High effort",
                font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Comparison table on the right
    tbl_x = Inches(6.8)
    tbl_y = Inches(1.4)

    headers = ["", "Manual", "ChatGPT", "Klue/Crayon", "MiAmar"]
    rows = [
        ["Freshness", "Days old", "Point-in-time", "Manual feed", "Daily"],
        ["Battlecards", "Copy-paste", "Hallucinations", "Semi-auto", "Auto + live"],
        ["Time to insight", "Hours", "Min (unreliable)", "Hours", "Seconds"],
        ["Cost", "Free", "$20/mo", "$20-45K/yr", "$199-499/mo"],
        ["Weekly effort", "6-8 hrs", "Per-query", "3-5 hrs", "~15 min"],
    ]

    col_widths = [Inches(1.1), Inches(0.95), Inches(1.1), Inches(1.15), Inches(1.15)]
    row_h = Inches(0.45)

    # Header row
    for c, (header, col_w) in enumerate(zip(headers, col_widths)):
        x = tbl_x + sum(w for w in [Inches(0)] + [col_widths[j] for j in range(c)])
        bg = BLUE if c == 4 else DARK if c > 0 else LIGHT_BG
        fg = WHITE if c >= 1 else DARK

        cell = add_rounded_rect(slide, x, tbl_y, col_w - Inches(0.02), row_h, bg)
        add_textbox(slide, x + Inches(0.05), tbl_y + Inches(0.05),
                    col_w - Inches(0.12), row_h - Inches(0.1),
                    header, font_size=10, color=fg, bold=True, alignment=PP_ALIGN.CENTER)

    # Data rows
    for r, row in enumerate(rows):
        y = tbl_y + (r + 1) * row_h
        for c, val in enumerate(row):
            x = tbl_x + sum(w for w in [Inches(0)] + [col_widths[j] for j in range(c)])
            bg = LIGHT_TEAL_BG if c == 4 else WHITE
            fg = TEAL if c == 4 else DARK if c == 0 else GRAY

            cell = add_rounded_rect(slide, x, y, col_widths[c] - Inches(0.02),
                                    row_h, bg)
            add_textbox(slide, x + Inches(0.05), y + Inches(0.05),
                        col_widths[c] - Inches(0.12), row_h - Inches(0.1),
                        val, font_size=9, color=fg, bold=(c == 0 or c == 4),
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "Here's how we position against the market. On one axis, you have the effort required — from fully automated to fully manual. On the other, the depth of intelligence output. The incumbents, Klue and Crayon, deliver depth but require significant analyst effort and $20K to $45K per year. Manual methods are free but take 6 to 8 hours per week and produce shallow, outdated intel. Generic AI tools like ChatGPT are fast but unreliable — they hallucinate competitive data. MiAmar sits in the top-left quadrant: deep intelligence, minimal effort, at a fraction of the incumbent cost. That's the position nobody occupies today.")


def slide_11_traction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Live Product, Early Traction, Clear Path",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # Left panel: Current Status
    left_x = Inches(0.8)
    left_w = Inches(5.5)

    left_card = add_rounded_rect(slide, left_x, Inches(1.4), left_w, Inches(5.5),
                                 LIGHT_TEAL_BG, TEAL)

    add_textbox(slide, left_x + Inches(0.3), Inches(1.6), left_w - Inches(0.6), Inches(0.5),
                "Today — Live & Working", font_size=20, color=TEAL, bold=True)

    status_items = [
        "Product is live and monitoring companies daily",
        "Full pipeline: LinkedIn posts, profiles, news, funding, job postings",
        "AI generates battlecards, event classifications, prospect cards, weekly newsletters",
        "Integrations: HubSpot, Slack, Monday.com",
        "27 qualified inbound leads identified from Reddit alone",
    ]

    txBox = slide.shapes.add_textbox(left_x + Inches(0.3), Inches(2.3),
                                     left_w - Inches(0.6), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(status_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.font.name = FONT_BODY
        p.space_after = Pt(12)

    # Right panel: Roadmap
    right_x = Inches(6.8)
    right_w = Inches(5.8)

    right_card = add_rounded_rect(slide, right_x, Inches(1.4), right_w, Inches(5.5),
                                  LIGHT_BG, BLUE)

    add_textbox(slide, right_x + Inches(0.3), Inches(1.6), right_w - Inches(0.6), Inches(0.5),
                "12-Month Roadmap", font_size=20, color=BLUE, bold=True)

    milestones = [
        ("Q1-Q2 2026", "Onboard first 10 paying customers.\nValidate pricing. Refine AI accuracy."),
        ("Q2-Q3 2026", "50 paying customers. Self-serve signup.\nSalesforce integration."),
        ("Q3-Q4 2026", "$50K MRR. Inbound marketing engine.\nExpand sources (Glassdoor, G2)."),
        ("Q1 2027", "Seed round position: $100K+ MRR,\n100+ customers, proven unit economics."),
    ]

    for i, (quarter, desc) in enumerate(milestones):
        y = Inches(2.3) + i * Inches(1.15)

        # Quarter badge
        badge = add_rounded_rect(slide, right_x + Inches(0.3), y,
                                 Inches(1.3), Inches(0.35), BLUE)
        add_textbox(slide, right_x + Inches(0.3), y,
                    Inches(1.3), Inches(0.35),
                    quarter, font_size=11, color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, right_x + Inches(1.8), y - Inches(0.05),
                    right_w - Inches(2.2), Inches(0.8),
                    desc, font_size=12, color=DARK)

    set_notes(slide,
              "We're not pitching a concept — the product is live. We're monitoring companies daily, generating battlecards and intelligence automatically, and delivering weekly digests. The AI pipeline works end to end. We've identified 27 high-intent inbound leads from Reddit communities alone — PMMs actively asking for exactly what we've built. Our 12-month plan is straightforward: onboard our first 10 paying customers in Q1-Q2, reach 50 customers by Q3, hit $50K MRR by end of year, and position for a seed round at $100K+ MRR with proven unit economics.")


def slide_12_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "Built by Operators, Not Observers",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # Eran - left column
    eran_x = Inches(0.8)
    card_w = Inches(5.5)
    card_h = Inches(3.8)

    eran_card = add_rounded_rect(slide, eran_x, Inches(1.4), card_w, card_h, LIGHT_BG, BLUE)

    add_textbox(slide, eran_x + Inches(0.3), Inches(1.6), card_w - Inches(0.6), Inches(0.45),
                "Eran Menachemi — CEO", font_size=22, color=BLUE, bold=True)

    eran_bullets = [
        "Product executive with deep GTM and product strategy experience",
        "Co-founder & COO at Navicula — built and scaled product & operations from zero",
        "Owns vision, product, and go-to-market at MiAmar",
        "Understands the PMM pain firsthand — lived it across multiple B2B companies",
    ]

    txBox = slide.shapes.add_textbox(eran_x + Inches(0.3), Inches(2.2),
                                     card_w - Inches(0.6), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(eran_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.font.name = FONT_BODY
        p.space_after = Pt(8)

    # Lior - right column
    lior_x = Inches(6.8)

    lior_card = add_rounded_rect(slide, lior_x, Inches(1.4), card_w, card_h, LIGHT_BG, TEAL)

    add_textbox(slide, lior_x + Inches(0.3), Inches(1.6), card_w - Inches(0.6), Inches(0.45),
                "Lior Cohen — CTO", font_size=22, color=TEAL, bold=True)

    lior_bullets = [
        "10 years spanning data engineering, full-stack development, and cybersecurity",
        "Built MiAmar's entire data pipeline: scraping infrastructure, AI classification, and product",
        "Designed for scale: automated monitoring of thousands of companies with daily refresh",
        "Technical depth to build fast and iterate without expensive engineering teams",
    ]

    txBox = slide.shapes.add_textbox(lior_x + Inches(0.3), Inches(2.2),
                                     card_w - Inches(0.6), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(lior_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.font.name = FONT_BODY
        p.space_after = Pt(8)

    # "Why this team" box at bottom
    why_y = Inches(5.5)
    why_card = add_rounded_rect(slide, Inches(0.8), why_y, Inches(11.7), Inches(1.6),
                                LIGHT_BLUE_BG, BLUE)

    add_textbox(slide, Inches(1.1), why_y + Inches(0.15), Inches(11.1), Inches(0.35),
                "Why This Team", font_size=16, color=BLUE, bold=True)

    why_lines = [
        "Eran knows the buyer — he's been the buyer. Product and GTM instincts from building and scaling a company.",
        "Lior knows the stack — built the entire product solo, from data infrastructure to AI to frontend.",
        "Lean two-person team that ships fast. No overhead. Every dollar goes into product and growth.",
    ]

    txBox = slide.shapes.add_textbox(Inches(1.1), why_y + Inches(0.5),
                                     Inches(11.1), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(why_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {line}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = FONT_BODY
        p.space_after = Pt(4)

    set_notes(slide,
              "We're a two-person founding team and that's by design. Eran is a product executive who co-founded Navicula, where he served as COO and built the product and operations from scratch. He understands the PMM buyer because he's been that buyer across multiple B2B companies. Lior has 10 years across data engineering, full-stack development, and cybersecurity. He built MiAmar's entire product solo — the scraping infrastructure, the AI classification pipeline, the frontend, the integrations. This team ships fast. We're not burning money on a large engineering team — we're two senior operators who can build, sell, and iterate rapidly. At the pre-seed stage, that's exactly what you want.")


def slide_13_the_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
                "$1M Pre-Seed to Capture the Mid-Market CI Gap",
                font_size=32, color=BLUE, bold=True, font_name=FONT_TITLE)

    # Big number
    hero = add_rounded_rect(slide, Inches(4.5), Inches(1.3), Inches(4.3), Inches(0.9),
                            BLUE)
    add_textbox(slide, Inches(4.5), Inches(1.35), Inches(4.3), Inches(0.8),
                "Raising: $1,000,000", font_size=30, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Use of funds — 3 cards
    funds = [
        ("Product & Engineering", "45% — $450K",
         "AI model improvements, new monitoring\nsources (Glassdoor, G2), Salesforce\nintegration, self-serve onboarding",
         BLUE),
        ("Sales & Marketing", "35% — $350K",
         "First sales hire, content marketing,\ncommunity presence (PMM communities),\npaid acquisition experiments",
         TEAL),
        ("Operations & Infra", "20% — $200K",
         "Cloud infrastructure for scaling,\ndata pipeline costs, legal,\n18-month runway buffer",
         ORANGE),
    ]

    card_w = Inches(3.8)
    gap = Inches(0.25)
    start_x = Inches(0.7)

    for i, (title, pct, desc, accent) in enumerate(funds):
        x = start_x + i * (card_w + gap)
        y = Inches(2.5)
        h = Inches(2.5)

        card = add_rounded_rect(slide, x, y, card_w, h, LIGHT_BG, accent)

        # Accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_textbox(slide, x + Inches(0.25), y + Inches(0.25), card_w - Inches(0.5), Inches(0.4),
                    title, font_size=16, color=accent, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.4),
                    pct, font_size=20, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.25), y + Inches(1.25), card_w - Inches(0.5), Inches(1.1),
                    desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 18-month milestones
    ms_y = Inches(5.3)
    ms_card = add_rounded_rect(slide, Inches(0.7), ms_y, Inches(11.9), Inches(1.7),
                               LIGHT_BLUE_BG, BLUE)

    add_textbox(slide, Inches(1.0), ms_y + Inches(0.15), Inches(11.3), Inches(0.4),
                "What $1M Gets Us To (18-Month Milestones)",
                font_size=16, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    milestones = [
        "100+ paying customers",
        "$100K+ MRR",
        "Proven unit economics\n(CAC, LTV, churn)",
        "Seed round ready\nwith real metrics",
    ]

    ms_card_w = Inches(2.6)
    ms_gap = Inches(0.2)
    ms_start_x = Inches(1.0)

    for i, ms in enumerate(milestones):
        x = ms_start_x + i * (ms_card_w + ms_gap)
        y = ms_y + Inches(0.6)

        badge = add_rounded_rect(slide, x, y, ms_card_w, Inches(0.85), WHITE, BLUE)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.08),
                    ms_card_w - Inches(0.2), Inches(0.7),
                    ms, font_size=13, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "We're raising $1 million to do three things: make the product best-in-class, build our initial customer base, and prove the unit economics for a seed round. Forty-five percent goes to product and engineering — improving AI accuracy, adding new data sources like Glassdoor and G2, building Salesforce integration, and enabling self-serve signup. Thirty-five percent goes to sales and marketing — our first sales hire, content marketing in PMM communities, and paid acquisition experiments. Twenty percent covers infrastructure and runway. This gets us to 100+ paying customers, $100K+ MRR, and proven unit economics within 18 months — a strong seed-round position.")


def slide_14_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Blue accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_textbox(slide, Inches(2), Inches(1.5), Inches(9.333), Inches(0.8),
                "Let's Talk", font_size=48, color=BLUE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # CTA box
    cta = add_rounded_rect(slide, Inches(4.0), Inches(2.8), Inches(5.333), Inches(0.8),
                           BLUE)
    add_textbox(slide, Inches(4.0), Inches(2.85), Inches(5.333), Inches(0.7),
                "Book a demo at miamar.io",
                font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Contact details
    contacts = [
        "Eran Menachemi — eran@miamar.io",
        "Lior Cohen — lior@miamar.io",
        "Website: miamar.io",
    ]

    for i, contact in enumerate(contacts):
        y = Inches(4.1) + i * Inches(0.5)
        add_textbox(slide, Inches(2), y, Inches(9.333), Inches(0.45),
                    contact, font_size=18, color=DARK, alignment=PP_ALIGN.CENTER)

    # Closing line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.5), Inches(5.6), Inches(2.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    add_textbox(slide, Inches(1.5), Inches(5.9), Inches(10.333), Inches(0.6),
                '"We\'re building the CI platform that PMMs have been asking for — literally. Come see it live."',
                font_size=16, color=GRAY, bold=False, alignment=PP_ALIGN.CENTER)

    set_notes(slide,
              "Thank you for your time. We'd love to show you the product live — it takes 60 seconds to add a competitor and see the AI work. Everything I showed you today is real, not mocked up. Reach out to either of us or book a demo directly at miamar.io. We're excited about this market and we'd love to have you on this journey with us.")


# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_market_pain(prs)
    slide_04_solution(prs)
    slide_05_how_it_works(prs)
    slide_06_features(prs)
    slide_07_market_size(prs)
    slide_08_why_now(prs)
    slide_09_business_model(prs)
    slide_10_competitive(prs)
    slide_11_traction(prs)
    slide_12_team(prs)
    slide_13_the_ask(prs)
    slide_14_contact(prs)

    out_path = "MiAmar-Pitch-Deck.pptx"
    prs.save(out_path)
    print(f"Deck saved to {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
